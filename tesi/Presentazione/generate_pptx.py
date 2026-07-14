from pptx import Presentation
from pptx.util import Inches, Pt

prs = Presentation()

# Layouts:
# 0: Title Slide
# 1: Title and Content

# Slide 1: Titolo
slide = prs.slides.add_slide(prs.slide_layouts[0])
slide.shapes.title.text = "Quenya: un traduttore da linguaggio naturale a query per Blazegraph"
slide.placeholders[1].text = "Alberto Zuccari\n\nRelatore: Chiar.mo Prof. Fabio Vitali\nCorrelatore: Chiar.mo Prof. Paolo Bonora\n\nUniversità di Bologna"

# Slide 2: Background 1 - Il dominio: DH.ARC e il Datavault
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Il dominio: DH.ARC e il Datavault"
tf = slide.placeholders[1].text_frame
tf.text = "Centro Ricerca Umanistica DH.ARC"
p = tf.add_paragraph()
p.text = "Gestisce il Datavault, un archivio di dati culturali."
p.level = 1
p = tf.add_paragraph()
p.text = "Si basa su un Knowledge Graph: ~19 Milioni di triple RDF."
p.level = 1
p = tf.add_paragraph()
p.text = "Il vantaggio dei dati a grafo: Collega tra loro dati culturali di natura molto diversa."
p.level = 0

# Slide 3: Background 2 - Il problema dell'accessibilità
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Il problema dell'accessibilità"
tf = slide.placeholders[1].text_frame
tf.text = "L'Utente: Ricercatore Umanistico"
p = tf.add_paragraph()
p.text = "Conosce i dati ma non sa programmare."
p.level = 1
p = tf.add_paragraph()
p.text = "Il problema pratico: La Barriera SPARQL"
p.level = 0
p = tf.add_paragraph()
p.text = "Per leggere e interrogare il Datavault serve saper scrivere in SPARQL, che ha una sintassi rigida e complessa."
p.level = 1

# Slide 4: Background 3 - I limiti delle soluzioni attuali
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "I limiti delle interfacce attuali"
tf = slide.placeholders[1].text_frame
tf.text = "I classici filtri non bastano"
p = tf.add_paragraph()
p.text = "Le interfacce a menu o a filtri limitano le vere potenzialità del grafo."
p.level = 1
p = tf.add_paragraph()
p.text = "Cosa serve davvero?"
p.level = 0
p = tf.add_paragraph()
p.text = "Un modo per fare ricerche complesse in libertà, scrivendo come si parla (Linguaggio Naturale)."
p.level = 1

# Slide 5: Quenya 1 - Cos'è e Obiettivi
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Obiettivo: Quenya"
tf = slide.placeholders[1].text_frame
tf.text = "Quenya: un traduttore da Domanda Naturale a Query Eseguibile"
p = tf.add_paragraph()
p.text = "Esempio: \"Quali manoscritti del XIV secolo sono digitalizzati?\" -> SELECT ?m WHERE { ... }"
p.level = 1
p = tf.add_paragraph()
p.text = "I tre pilastri del progetto:"
p.level = 0
p = tf.add_paragraph()
p.text = "1. Affidabilità: Capire davvero cosa sta cercando l'utente."
p.level = 1
p = tf.add_paragraph()
p.text = "2. Nessun errore: Bloccare le allucinazioni del modello prima che arrivino al database."
p.level = 1
p = tf.add_paragraph()
p.text = "3. Esperienza d'uso: Una semplice barra di ricerca, senza complicazioni."
p.level = 1

# Slide 6: Quenya 2 - Architettura del Sistema
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Architettura del Sistema"
tf = slide.placeholders[1].text_frame
tf.text = "Frontend"
p = tf.add_paragraph()
p.text = "Semplice barra di ricerca per l'inserimento della query."
p.level = 1
p = tf.add_paragraph()
p.text = "Orchestrator (Backend)"
p.level = 0
p = tf.add_paragraph()
p.text = "Controllo e gestione del flusso tra utente, LLM e Database."
p.level = 1
p = tf.add_paragraph()
p.text = "Modello LLM (Fine-tuned)"
p.level = 0
p = tf.add_paragraph()
p.text = "Riceve il prompt e genera la query SPARQL."
p.level = 1
p = tf.add_paragraph()
p.text = "Blazegraph"
p.level = 0
p = tf.add_paragraph()
p.text = "Il database RDF che esegue la query e restituisce i risultati."
p.level = 1

# Slide 7: Quenya 3 - Il Motore LLM (LoRA) e Gestione Allucinazioni
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Il Motore LLM e Gestione Allucinazioni"
tf = slide.placeholders[1].text_frame
tf.text = "La scelta: LoRA Fine-tuning"
p = tf.add_paragraph()
p.text = "I grandi modelli commerciali sono lenti e non conoscono il nostro database specifico."
p.level = 1
p = tf.add_paragraph()
p.text = "Quenya è un modello locale e leggero, addestrato esattamente sui dati del Datavault."
p.level = 1
p = tf.add_paragraph()
p.text = "Risolvere il problema delle allucinazioni"
p.level = 0
p = tf.add_paragraph()
p.text = "Rischio: Il modello inventa proprietà che non esistono."
p.level = 1
p = tf.add_paragraph()
p.text = "Soluzione: Generazione Vincolata (forziamo l'IA a pescare solo dal vocabolario ufficiale)."
p.level = 1
p = tf.add_paragraph()
p.text = "Risultato: Le query generate non vanno in errore di sintassi."
p.level = 1

# Slide 8: Valutazione 1 - Metodologia e Dataset
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Valutazione: Metodologia e Dataset"
tf = slide.placeholders[1].text_frame
tf.text = "I Dati di Test"
p = tf.add_paragraph()
p.text = "Training Set: Visti in fase di addestramento."
p.level = 1
p = tf.add_paragraph()
p.text = "Test Set: Dati tenuti nascosti per la validazione."
p.level = 1
p = tf.add_paragraph()
p.text = "Cosa abbiamo misurato"
p.level = 0
p = tf.add_paragraph()
p.text = "Efficacia: Precision, Recall e F1-Score."
p.level = 1
p = tf.add_paragraph()
p.text = "Efficienza: Percentuale di query eseguite con successo e tempi di attesa."
p.level = 1
p = tf.add_paragraph()
p.text = "Tipi di query testate: Ricerche semplici, combinazioni di filtri, relazioni complesse."
p.level = 1

# Slide 9: Valutazione 2 - Risultati
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Valutazione: Risultati Ottenuti"
tf = slide.placeholders[1].text_frame
tf.text = "Tre grandi successi:"
p = tf.add_paragraph()
p.text = "Alta Accuratezza nell'interpretazione della domanda."
p.level = 1
p = tf.add_paragraph()
p.text = "Bassa Latenza: Tempi di risposta naturali per l'utente web."
p.level = 1
p = tf.add_paragraph()
p.text = "Zero Errori Sintattici grazie alla Generazione Vincolata."
p.level = 1
p = tf.add_paragraph()
p.text = "Considerazioni finali:"
p.level = 0
p = tf.add_paragraph()
p.text = "LLM e Validatore si districano bene anche in relazioni RDF frammentate."
p.level = 1
p = tf.add_paragraph()
p.text = "Il modello compatto regge il confronto con i giganti commerciali su questo dominio."
p.level = 1

# Slide 10: Conclusioni e Sviluppi Futuri
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Conclusioni e Sviluppi Futuri"
tf = slide.placeholders[1].text_frame
tf.text = "Conclusioni"
p = tf.add_paragraph()
p.text = "Dati per tutti: Nessun bisogno di imparare SPARQL per i ricercatori."
p.level = 1
p = tf.add_paragraph()
p.text = "Efficienza su misura: Ottimi risultati senza supercomputer."
p.level = 1
p = tf.add_paragraph()
p.text = "Il patrimonio del DH.ARC è interrogabile come un motore di ricerca."
p.level = 1
p = tf.add_paragraph()
p.text = "Sviluppi Futuri"
p.level = 0
p = tf.add_paragraph()
p.text = "Chat: Più turni di dialogo per chiarire query ambigue."
p.level = 1
p = tf.add_paragraph()
p.text = "Lingue: Supporto per l'Inglese."
p.level = 1
p = tf.add_paragraph()
p.text = "Complessità: Federated SPARQL e Query geospaziali."
p.level = 1

prs.save("presentazione_tesi_10_slide.pptx")
print("Presentation saved successfully as presentazione_tesi_10_slide.pptx")
